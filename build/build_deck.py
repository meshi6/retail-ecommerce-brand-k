"""Build the Kiabi x Grupo One B2B reporting template (.pptx).

Brand sources:
- Colors & typography: assets/kia_brandbook_charte_eng.pdf (p.20 colours, p.58 Figtree)
- Kiabi logo: assets/kiabi-logo.png (navy, light bg) / build/kiabi_white.png (dark bg)
- Grupo One logo: build/gone.png / build/gone_white.png
- Motifs: build/make_motifs.py (Art Deco sunburst - Paris / Buenos Aires bridge)
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
NAVY = RGBColor(0x04, 0x00, 0x37)      # Kiabi blue (main)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x00, 0x6E, 0xFB)      # secondary blue
ORANGE = RGBColor(0xFB, 0x5B, 0x00)
RED = RGBColor(0xE9, 0x1B, 0x2E)
YELLOW = RGBColor(0xFF, 0xDA, 0x00)
GREEN = RGBColor(0x20, 0xB1, 0x4A)


def tint(c, f):
    """f = fraction of colour kept; rest white (brand book allows 70/40/10% tints)."""
    return RGBColor(*(int(ch * f + 255 * (1 - f)) for ch in (c[0], c[1], c[2])))


def shade_on_navy(c, f):
    return RGBColor(*(int(ch * f + n * (1 - f)) for ch, n in zip((c[0], c[1], c[2]), (4, 0, 55))))


BODY = tint(NAVY, 0.70)        # body text
MUTED = tint(NAVY, 0.45)       # captions
HAIR = tint(NAVY, 0.15)        # hairlines
LAVENDER = RGBColor(0xBD, 0xBB, 0xD4)  # light text on navy
FONT = "Figtree"

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]

ACCENTS = {"blue": BLUE, "orange": ORANGE, "red": RED, "green": GREEN,
           "yellow": YELLOW, "navy": NAVY}


# ---------------------------------------------------------------- helpers
def slide_new(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY if dark else WHITE
    return s


def _set_spc(run, val):
    run._r.get_or_add_rPr().set("spc", str(val))


def text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts {runs:[(txt,{size,color,bold,italic,font,spc})],
    align, space_after(pt), line(float)}"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in (tf, ):
        pass
    tb.text_frame.margin_left = tb.text_frame.margin_right = 0
    tb.text_frame.margin_top = tb.text_frame.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if "space_after" in p:
            para.space_after = Pt(p["space_after"])
        if "space_before" in p:
            para.space_before = Pt(p["space_before"])
        if "line" in p:
            para.line_spacing = p["line"]
        for t, st in p["runs"]:
            r = para.add_run()
            r.text = t
            f = r.font
            f.name = st.get("font", FONT)
            f.size = Pt(st.get("size", 12))
            f.color.rgb = st.get("color", NAVY)
            f.bold = st.get("bold", False)
            f.italic = st.get("italic", False)
            if "spc" in st:
                _set_spc(r, st["spc"])
    return tb


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, rounded=False,
         radius=0.06, dashed=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
        if dashed:
            ln = shp.line._get_or_add_ln()
            d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
            ln.append(d)
    shp.shadow.inherit = False
    return shp


def circle(slide, cx, cy, d, fill=None, line=None, line_w=1.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2),
                                 Inches(cy - d / 2), Inches(d), Inches(d))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def hline(slide, x, y, w, color, weight=1.0):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def ribbon(slide, y=SH - 0.14, h=0.14):
    """Five-colour brand strip (secondary palette, brand book p.20)."""
    seg = SW / 5
    for i, c in enumerate([ORANGE, RED, YELLOW, GREEN, BLUE]):
        rect(slide, i * seg, y, seg + 0.01, h, fill=c)


def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt


def footer(slide, idx):
    hline(slide, 0.55, 7.08, SW - 1.1, HAIR, 1.0)
    text(slide, 0.55, 7.16, 6.0, 0.3, [{"runs": [
        ("GRUPO ONE × KIABI", {"size": 7.5, "color": MUTED, "bold": True, "spc": 200}),
        ("   ·   CONFIDENTIAL — FOR KIABI HQ", {"size": 7.5, "color": MUTED, "spc": 200}),
    ]}])
    text(slide, SW - 4.55, 7.16, 4.0, 0.3, [{"runs": [
        ("[MONTH YEAR]   ·   ", {"size": 7.5, "color": MUTED, "spc": 200}),
        (f"{idx:02d} / 10", {"size": 7.5, "color": MUTED, "bold": True, "spc": 200}),
    ], "align": PP_ALIGN.RIGHT}])


def header(slide, num, eyebrow, title, accent_name):
    accent = ACCENTS[accent_name]
    slide.shapes.add_picture("build/fan_light.png", Inches(SW - 2.1), 0,
                             Inches(2.1), Inches(2.1))
    slide.shapes.add_picture(f"build/glyph_{accent_name}.png",
                             Inches(0.55), Inches(0.56), Inches(0.30), Inches(0.30))
    text(slide, 0.97, 0.47, 10.0, 0.3, [{"runs": [
        (f"{num:02d} — {eyebrow}", {"size": 10.5, "color": accent, "bold": True, "spc": 280}),
    ]}])
    text(slide, 0.95, 0.74, 10.5, 0.65, [{"runs": [
        (title, {"size": 27, "color": NAVY, "bold": True}),
    ]}])
    hline(slide, 0.55, 1.52, SW - 1.1, HAIR, 1.0)
    footer(slide, num)
    return accent


def chip(slide, x, y, w, h, label, fill=None, line=None, txt_color=WHITE,
         size=9, bold=True, spc=150):
    shp = rect(slide, x, y, w, h, fill=fill, line=line, line_w=1.0, rounded=True,
               radius=0.5)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = txt_color
    if spc:
        _set_spc(r, spc)
    return shp


# ================================================================ 1 · COVER
s = slide_new(dark=True)
s.shapes.add_picture("build/fan_dark.png", Inches(SW - 5.3), Inches(SH - 5.3),
                     Inches(5.3), Inches(5.3))
# logos
s.shapes.add_picture("build/kiabi_white.png", Inches(0.55), Inches(0.52),
                     Inches(1.95), Inches(1.95 * 320 / 1440))
gw = 1.05
s.shapes.add_picture("build/gone_white.png", Inches(SW - 0.55 - gw), Inches(0.42),
                     Inches(gw), Inches(gw * 908 / 900))
# eyebrow + title
text(s, 0.57, 2.42, 9.5, 0.35, [{"runs": [
    ("PARIS · BUENOS AIRES", {"size": 12, "color": YELLOW, "bold": True, "spc": 380}),
    ("   —   BRAND PERFORMANCE REPORT", {"size": 12, "color": LAVENDER, "spc": 380}),
]}])
text(s, 0.55, 2.85, 9.0, 1.7, [
    {"runs": [("[REPORT TITLE]", {"size": 52, "color": WHITE, "bold": True})], "line": 1.0},
    {"runs": [("Kiabi Argentina — progress, performance & insights",
               {"size": 17, "color": LAVENDER})], "space_before": 14},
])
# meta chips
cy = 5.30
chip(s, 0.55, cy, 1.95, 0.46, "EDITION Nº [00]", line=LAVENDER,
     txt_color=WHITE, size=10, spc=200)
chip(s, 2.70, cy, 1.85, 0.46, "[MONTH YEAR]", line=LAVENDER, txt_color=WHITE,
     size=10, spc=200)
chip(s, 4.75, cy, 3.30, 0.46, "PREPARED FOR KIABI HQ · FRANCE", line=LAVENDER,
     txt_color=WHITE, size=10, spc=200)
# presenter line
text(s, 0.55, 6.55, 8.5, 0.5, [{"runs": [
    ("Presented by Grupo One — ", {"size": 10.5, "color": LAVENDER, "bold": True}),
    ("elevating global brands in the region", {"size": 10.5, "color": LAVENDER, "italic": True}),
]}])
ribbon(s)
notes(s, "PROMPTS — 1) Open with the single most important headline of this edition "
         "(one sentence). 2) Confirm the reporting period and what has changed since "
         "the last edition before moving on.")

# ====================================================== 2 · EXECUTIVE SUMMARY
s = slide_new()
acc = header(s, 2, "SYNTHÈSE · RESUMEN EJECUTIVO", "Executive Summary", "blue")
card_acc = [BLUE, ORANGE, GREEN, YELLOW]
labels = ["[KPI 1 — e.g. NET SALES]", "[KPI 2 — e.g. STORE TRAFFIC]",
          "[KPI 3 — e.g. CONVERSION]", "[KPI 4 — e.g. SOCIAL REACH]"]
for i in range(4):
    x = 0.55 + i * 3.1075
    w = 2.92
    rect(s, x, 1.85, w, 2.05, fill=tint(card_acc[i], 0.08), rounded=True, radius=0.07)
    rect(s, x, 1.85, w, 0.09, fill=card_acc[i], rounded=False)
    text(s, x + 0.22, 2.18, w - 0.44, 0.8, [{"runs": [
        ("[METRIC]", {"size": 34, "color": NAVY, "bold": True})]}])
    text(s, x + 0.22, 2.95, w - 0.44, 0.45, [{"runs": [
        (labels[i], {"size": 9, "color": BODY, "bold": True, "spc": 120})], "line": 1.1}])
    text(s, x + 0.22, 3.42, w - 0.44, 0.35, [{"runs": [
        ("▲ [+X.X%] ", {"size": 9.5, "color": GREEN, "bold": True}),
        ("VS LAST MONTH", {"size": 8, "color": MUTED, "spc": 100}),
    ]}])
# headline insight band
rect(s, 0.55, 4.25, SW - 1.1, 1.65, fill=NAVY, rounded=True, radius=0.09)
s.shapes.add_picture("build/glyph_yellow.png", Inches(0.92), Inches(4.62),
                     Inches(0.34), Inches(0.34))
text(s, 1.45, 4.50, 11.0, 0.3, [{"runs": [
    ("HEADLINE INSIGHT", {"size": 10, "color": YELLOW, "bold": True, "spc": 300})]}])
text(s, 1.45, 4.82, 10.95, 0.95, [{"runs": [
    ("[INSIGHT — the one sentence Kiabi HQ should remember from this report.]",
     {"size": 16.5, "color": WHITE})], "line": 1.15}])
text(s, 0.55, 6.12, 12.2, 0.4, [{"runs": [
    ("Context — ", {"size": 10.5, "color": NAVY, "bold": True}),
    ("[Optional supporting fact, proof point or caveat for the headline above.]",
     {"size": 10.5, "color": BODY}),
]}])
notes(s, "PROMPTS — 1) Walk the four KPIs left to right; flag any figure that "
         "deviates >10% from plan. 2) State the headline insight in your own words "
         "and link it to the decision you want from HQ.")

# ======================================================== 3 · MARKET SNAPSHOT
s = slide_new()
acc = header(s, 3, "PANORAMA DU MARCHÉ · PANORAMA DEL MERCADO", "Market Snapshot", "green")
text(s, 0.55, 1.82, 6.3, 0.4, [{"runs": [
    ("Argentina retail & fashion context", {"size": 14.5, "color": NAVY, "bold": True})]}])
bullets = [
    "[CONTEXT 1 — macro: inflation, FX, consumer purchasing power.]",
    "[CONTEXT 2 — retail: traffic, mall performance, promo intensity.]",
    "[CONTEXT 3 — fashion: category trends, competitor moves.]",
    "[CONTEXT 4 — regulation / imports / pricing environment.]",
]
text(s, 0.55, 2.35, 6.3, 2.6,
     [{"runs": [("—  ", {"size": 12, "color": GREEN, "bold": True}),
                (b, {"size": 12, "color": BODY})],
       "space_after": 12, "line": 1.15} for b in bullets])
# watch-this callout
rect(s, 0.55, 5.05, 6.3, 1.30, fill=tint(YELLOW, 0.16), rounded=True, radius=0.10)
rect(s, 0.55, 5.05, 0.09, 1.30, fill=YELLOW)
text(s, 0.85, 5.22, 5.8, 0.3, [{"runs": [
    ("WATCH THIS", {"size": 9.5, "color": NAVY, "bold": True, "spc": 300})]}])
text(s, 0.85, 5.52, 5.75, 0.75, [{"runs": [
    ("[MARKET SIGNAL — one emerging trend or risk HQ should keep on its radar.]",
     {"size": 11.5, "color": BODY})], "line": 1.15}])
# right data tiles
tile_acc = [GREEN, BLUE, ORANGE]
tile_lbl = ["[INDICATOR 1 — e.g. APPAREL MARKET YoY]",
            "[INDICATOR 2 — e.g. CPI / FASHION INFLATION]",
            "[INDICATOR 3 — e.g. E-COMMERCE SHARE]"]
for i in range(3):
    y = 1.85 + i * 1.55
    rect(s, 7.30, y, 5.48, 1.38, fill=tint(tile_acc[i], 0.08), rounded=True, radius=0.09)
    rect(s, 7.30, y, 0.09, 1.38, fill=tile_acc[i])
    text(s, 7.62, y + 0.18, 1.9, 0.7, [{"runs": [
        ("[X.X%]", {"size": 26, "color": NAVY, "bold": True})]}])
    text(s, 9.55, y + 0.24, 3.1, 0.62, [{"runs": [
        (tile_lbl[i], {"size": 9.5, "color": BODY, "bold": True})], "line": 1.1}])
    text(s, 9.55, y + 0.92, 3.1, 0.3, [{"runs": [
        ("SOURCE: [INDEC / CAME / KANTAR] · [DATE]", {"size": 7.5, "color": MUTED, "spc": 80})]}])
notes(s, "PROMPTS — 1) Give HQ the 30-second read on Argentina this month — what "
         "changed in the operating environment. 2) Translate one macro data point "
         "into a concrete implication for Kiabi pricing or assortment.")

# ================================================= 4 · BRAND LAUNCH PROGRESS
s = slide_new()
acc = header(s, 4, "LANCEMENT · AVANCE DEL LANZAMIENTO", "Brand Launch Progress", "orange")
# legend
text(s, 8.2, 1.78, 4.58, 0.3, [{"runs": [
    ("●", {"size": 11, "color": GREEN}), (" COMPLETED    ", {"size": 8.5, "color": BODY, "spc": 100}),
    ("●", {"size": 11, "color": YELLOW}), (" IN PROGRESS    ", {"size": 8.5, "color": BODY, "spc": 100}),
    ("●", {"size": 11, "color": HAIR}), (" UPCOMING", {"size": 8.5, "color": BODY, "spc": 100}),
], "align": PP_ALIGN.RIGHT}])
AXIS_Y = 3.95
hline(s, 0.9, AXIS_Y, 11.5, HAIR, 2.5)
hline(s, 0.9, AXIS_Y, 5.3, GREEN, 2.5)   # progress portion up to current node
node_x = [1.2, 3.7, 6.2, 8.7, 11.2]
status = ["done", "done", "now", "next", "next"]
ms = ["[MILESTONE 1 — e.g. Lease signed]", "[MILESTONE 2 — e.g. First shipment]",
      "[MILESTONE 3 — e.g. Store fit-out]", "[MILESTONE 4 — e.g. Soft opening]",
      "[MILESTONE 5 — e.g. Grand opening]"]
for i, x in enumerate(node_x):
    st = status[i]
    if st == "done":
        circle(s, x, AXIS_Y, 0.30, fill=GREEN)
        circle(s, x, AXIS_Y, 0.12, fill=WHITE)
    elif st == "now":
        circle(s, x, AXIS_Y, 0.38, fill=WHITE, line=YELLOW, line_w=3.0)
        circle(s, x, AXIS_Y, 0.14, fill=YELLOW)
    else:
        circle(s, x, AXIS_Y, 0.30, fill=WHITE, line=HAIR, line_w=2.0)
    above = i % 2 == 0
    ly = 2.55 if above else 4.35
    bx = x - 1.05
    text(s, bx, ly, 2.1, 0.55, [{"runs": [
        (ms[i], {"size": 10.5, "color": NAVY, "bold": True})],
        "align": PP_ALIGN.CENTER, "line": 1.05}])
    text(s, bx, ly + 0.62, 2.1, 0.3, [{"runs": [
        ("[DATE]", {"size": 9.5, "color": MUTED})], "align": PP_ALIGN.CENTER}])
# current phase detail
rect(s, 0.55, 5.62, SW - 1.1, 1.22, fill=tint(ORANGE, 0.08), rounded=True, radius=0.10)
rect(s, 0.55, 5.62, 0.09, 1.22, fill=ORANGE)
text(s, 0.88, 5.78, 2.2, 0.3, [{"runs": [
    ("CURRENT PHASE", {"size": 9.5, "color": ORANGE, "bold": True, "spc": 250})]}])
text(s, 0.88, 6.08, 11.6, 0.65, [{"runs": [
    ("[DETAIL — status of the active milestone: % complete, blockers, date risk, "
     "what HQ support is needed.]", {"size": 11.5, "color": BODY})], "line": 1.15}])
notes(s, "PROMPTS — 1) Confirm which milestones moved since last edition and "
         "whether the critical path date still holds. 2) Name the single biggest "
         "blocker on the current phase and the help needed from HQ, if any.")

# ==================================================== 5 · SALES & PERFORMANCE
s = slide_new()
acc = header(s, 5, "VENTES & PERFORMANCE · VENTAS Y DESEMPEÑO", "Sales & Performance", "blue")
text(s, 0.55, 1.80, 7.5, 0.35, [{"runs": [
    ("[METRIC — e.g. Net sales, ARS m] ", {"size": 12.5, "color": NAVY, "bold": True}),
    ("· actual vs target, last 6 months", {"size": 11, "color": MUTED}),
]}])
cd = CategoryChartData()
cd.categories = ["[M1]", "[M2]", "[M3]", "[M4]", "[M5]", "[M6]"]
cd.add_series("[ACTUAL]", (40, 55, 48, 62, 70, 78))
cd.add_series("[TARGET]", (45, 50, 55, 60, 66, 72))
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55),
                            Inches(2.25), Inches(7.55), Inches(4.55), cd)
ch = gframe.chart
ch.has_title = False
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(9)
ch.legend.font.name = FONT
ch.legend.font.color.rgb = BODY
plot = ch.plots[0]
plot.gap_width = 90
plot.overlap = -20
for ser, col in zip(plot.series, [BLUE, tint(NAVY, 0.25)]):
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = col
    ser.format.line.fill.background()
cax, vax = ch.category_axis, ch.value_axis
for ax in (cax, vax):
    ax.tick_labels.font.size = Pt(9)
    ax.tick_labels.font.name = FONT
    ax.tick_labels.font.color.rgb = MUTED
    ax.format.line.color.rgb = HAIR
cax.has_major_gridlines = False
vax.has_major_gridlines = True
vax.major_gridlines.format.line.color.rgb = tint(NAVY, 0.08)
vax.major_gridlines.format.line.width = Pt(0.75)
# commentary panel
rect(s, 8.45, 1.85, 4.33, 4.98, fill=tint(BLUE, 0.07), rounded=True, radius=0.06)
text(s, 8.75, 2.08, 3.8, 0.3, [{"runs": [
    ("COMMENTARY", {"size": 10, "color": BLUE, "bold": True, "spc": 300})]}])
blocks = [
    ("▲  WHAT WORKED", GREEN, "[INSIGHT — driver of over-performance: category, store, campaign.]"),
    ("▼  WHAT LAGGED", RED, "[INSIGHT — where we missed and the root cause.]"),
    ("→  ACTION / ASK", ORANGE, "[ACTION — corrective step taken, or decision requested from HQ.]"),
]
for i, (h, c, b) in enumerate(blocks):
    y = 2.55 + i * 1.42
    text(s, 8.75, y, 3.75, 0.3, [{"runs": [
        (h, {"size": 9.5, "color": c, "bold": True, "spc": 150})]}])
    text(s, 8.75, y + 0.30, 3.75, 0.95, [{"runs": [
        (b, {"size": 10.5, "color": BODY})], "line": 1.12}])
notes(s, "PROMPTS — 1) Explain the gap between actual and target — is it traffic, "
         "conversion or basket? 2) Close with the one corrective action and the "
         "date by which HQ will see its effect.")

# ====================================================== 6 · CONSUMER INSIGHTS
s = slide_new()
acc = header(s, 6, "INSIGHTS CONSOMMATEURS · INSIGHTS DEL CONSUMIDOR", "Consumer Insights", "red")
text(s, 0.55, 1.80, 5.9, 0.3, [{"runs": [
    ("BY THE NUMBERS", {"size": 10, "color": RED, "bold": True, "spc": 300})]}])
stats = [
    ("[X%]", "[QUANT 1 — e.g. brand awareness in Buenos Aires]"),
    ("[X.X]", "[QUANT 2 — e.g. NPS / satisfaction score]"),
    ("[X%]", "[QUANT 3 — e.g. repeat purchase intent]"),
]
for i, (num, lbl) in enumerate(stats):
    y = 2.25 + i * 1.18
    text(s, 0.55, y, 1.75, 0.7, [{"runs": [
        (num, {"size": 30, "color": NAVY, "bold": True})]}])
    text(s, 2.45, y + 0.10, 3.95, 0.65, [{"runs": [
        (lbl, {"size": 11, "color": BODY})], "line": 1.12}])
    if i < 2:
        hline(s, 0.55, y + 1.02, 5.85, HAIR, 0.75)
text(s, 0.55, 5.85, 5.85, 0.35, [{"runs": [
    ("SOURCE: [STUDY / PANEL / SAMPLE SIZE] · [DATE]", {"size": 8, "color": MUTED, "spc": 100})]}])
# right: qualitative
text(s, 6.95, 1.80, 5.85, 0.3, [{"runs": [
    ("VOICE OF THE CUSTOMER", {"size": 10, "color": RED, "bold": True, "spc": 300})]}])
for i in range(2):
    y = 2.25 + i * 1.80
    rect(s, 6.95, y, 5.83, 1.62, fill=tint(RED, 0.06), rounded=True, radius=0.09)
    rect(s, 6.95, y, 0.09, 1.62, fill=RED)
    text(s, 7.30, y + 0.20, 5.25, 0.85, [{"runs": [
        ("“[QUOTE — verbatim from a customer, store team or social comment.]”",
         {"size": 12, "color": NAVY, "italic": True})], "line": 1.15}])
    text(s, 7.30, y + 1.18, 5.25, 0.3, [{"runs": [
        ("— [SOURCE: NPS VERBATIM / STORE VISIT / IG COMMENT] · [DATE]",
         {"size": 8, "color": MUTED, "spc": 100})]}])
# takeaway
s.shapes.add_picture("build/glyph_red.png", Inches(0.55), Inches(6.32),
                     Inches(0.26), Inches(0.26))
text(s, 0.95, 6.28, 11.8, 0.55, [{"runs": [
    ("So what — ", {"size": 12, "color": NAVY, "bold": True}),
    ("[IMPLICATION for assortment, pricing or communication next month.]",
     {"size": 12, "color": BODY}),
]}])
notes(s, "PROMPTS — 1) Pair one number with one quote to tell a single customer "
         "story, not two lists. 2) End on the 'so what' — what this changes in "
         "next month's plan.")

# ======================================================== 7 · SOCIAL & DIGITAL
s = slide_new()
acc = header(s, 7, "SOCIAL & DIGITAL · REDES Y DIGITAL", "Social & Digital", "blue")
plats = [("INSTAGRAM", BLUE), ("TIKTOK", NAVY), ("FACEBOOK", GREEN), ("WHATSAPP / CRM", ORANGE)]
for i, (p, c) in enumerate(plats):
    x = 0.55 + i * 3.1075
    w = 2.92
    rect(s, x, 1.82, w, 1.62, fill=WHITE, line=HAIR, line_w=1.0, rounded=True, radius=0.09)
    rect(s, x, 1.82, w, 0.07, fill=c)
    text(s, x + 0.2, 2.02, w - 0.4, 0.28, [{"runs": [
        (p, {"size": 9.5, "color": MUTED, "bold": True, "spc": 200})]}])
    text(s, x + 0.2, 2.32, w - 0.4, 0.55, [{"runs": [
        ("[METRIC]", {"size": 24, "color": NAVY, "bold": True})]}])
    text(s, x + 0.2, 2.92, w - 0.4, 0.35, [{"runs": [
        ("▲ [+X%]  ", {"size": 9.5, "color": GREEN, "bold": True}),
        ("ENG. [X.X%]", {"size": 9, "color": BODY, "spc": 80}),
    ]}])
# engagement breakdown bars
text(s, 0.55, 3.72, 7.5, 0.3, [{"runs": [
    ("SHARE OF ENGAGEMENT BY PLATFORM", {"size": 10, "color": BLUE, "bold": True, "spc": 250})]}])
bar_data = [("INSTAGRAM", 0.80, BLUE), ("TIKTOK", 0.58, NAVY),
            ("FACEBOOK", 0.36, GREEN), ("WHATSAPP / CRM", 0.22, ORANGE)]
BARS_X, BARS_W = 2.45, 4.7
for i, (lbl, frac, c) in enumerate(bar_data):
    y = 4.18 + i * 0.64
    text(s, 0.55, y + 0.02, 1.8, 0.3, [{"runs": [
        (lbl, {"size": 8.5, "color": BODY, "bold": True, "spc": 80})]}])
    rect(s, BARS_X, y, BARS_W, 0.30, fill=tint(NAVY, 0.05), rounded=True, radius=0.5)
    rect(s, BARS_X, y, BARS_W * frac, 0.30, fill=c, rounded=True, radius=0.5)
    text(s, BARS_X + BARS_W * frac + 0.12, y + 0.02, 0.9, 0.3, [{"runs": [
        ("[X%]", {"size": 9.5, "color": NAVY, "bold": True})]}])
# highlight panel
rect(s, 8.45, 3.72, 4.33, 3.10, fill=tint(BLUE, 0.07), rounded=True, radius=0.07)
text(s, 8.73, 3.93, 3.8, 0.3, [{"runs": [
    ("CONTENT HIGHLIGHT", {"size": 9.5, "color": BLUE, "bold": True, "spc": 250})]}])
ph = rect(s, 8.73, 4.28, 1.95, 1.95, fill=WHITE, line=tint(BLUE, 0.45),
          line_w=1.25, rounded=True, radius=0.06, dashed=True)
text(s, 8.73, 4.95, 1.95, 0.6, [
    {"runs": [("DROP CREATIVE", {"size": 8, "color": tint(BLUE, 0.6), "bold": True, "spc": 100})],
     "align": PP_ALIGN.CENTER, "line": 1.1},
    {"runs": [("HERE", {"size": 8, "color": tint(BLUE, 0.6), "bold": True, "spc": 100})],
     "align": PP_ALIGN.CENTER, "line": 1.1},
])
text(s, 10.85, 4.32, 1.75, 1.4, [{"runs": [
    ("[BEST POST — reach, saves, sentiment]", {"size": 9.5, "color": BODY})], "line": 1.15}])
text(s, 8.73, 6.32, 3.8, 0.4, [{"runs": [
    ("[LINK OR CAMPAIGN NAME]", {"size": 8.5, "color": BLUE, "bold": True, "spc": 80})]}])
notes(s, "PROMPTS — 1) Lead with the channel that over- or under-delivered, not "
         "the full grid. 2) Show the highlighted creative and say why it worked — "
         "and whether it is replicable for the next campaign.")

# ====================================================== 8 · CHALLENGES & RISKS
s = slide_new()
acc = header(s, 8, "DÉFIS & RISQUES · DESAFÍOS Y RIESGOS", "Challenges & Risks", "red")
sev = [("HIGH", RED, WHITE), ("MEDIUM", YELLOW, NAVY), ("LOW", GREEN, WHITE)]
for i in range(3):
    y = 1.85 + i * 1.66
    label, c, tc = sev[i]
    # risk card
    rect(s, 0.55, y, 6.85, 1.50, fill=WHITE, line=HAIR, line_w=1.0, rounded=True, radius=0.08)
    rect(s, 0.55, y, 0.09, 1.50, fill=c)
    chip(s, 0.88, y + 0.22, 1.05, 0.34, label, fill=c, txt_color=tc, size=8.5, spc=150)
    text(s, 2.12, y + 0.22, 5.1, 0.35, [{"runs": [
        (f"[RISK {i + 1} — name the challenge honestly]",
         {"size": 12, "color": NAVY, "bold": True})]}])
    text(s, 0.88, y + 0.68, 6.25, 0.72, [{"runs": [
        ("[DESCRIPTION — impact on sales, opening date or brand, and likelihood.]",
         {"size": 10.5, "color": BODY})], "line": 1.12}])
    # mitigation card
    rect(s, 7.62, y, 5.16, 1.50, fill=tint(GREEN, 0.07), rounded=True, radius=0.08)
    text(s, 7.92, y + 0.18, 4.6, 0.28, [{"runs": [
        ("→  MITIGATION", {"size": 9, "color": GREEN, "bold": True, "spc": 200})]}])
    text(s, 7.92, y + 0.50, 4.6, 0.62, [{"runs": [
        ("[ACTION — what we are doing about it]", {"size": 10.5, "color": BODY})],
        "line": 1.12}])
    text(s, 7.92, y + 1.14, 4.6, 0.28, [{"runs": [
        ("OWNER: [NAME] · DUE: [DATE]", {"size": 8, "color": MUTED, "spc": 100})]}])
notes(s, "PROMPTS — 1) Be candid: present the highest-severity risk first and "
         "say plainly if it threatens a committed date or number. 2) For each risk, "
         "confirm the owner and whether mitigation is on track since last edition.")

# ============================================== 9 · NEXT STEPS & ACTION ITEMS
s = slide_new()
acc = header(s, 9, "PROCHAINES ÉTAPES · PRÓXIMOS PASOS", "Next Steps & Action Items", "green")
rows, cols = 6, 5
gframe = s.shapes.add_table(rows, cols, Inches(0.55), Inches(1.85),
                            Inches(12.23), Inches(3.55))
tbl = gframe.table
tbl.first_row = False
tbl.horz_banding = False
widths = [0.55, 5.65, 2.20, 1.90, 1.93]
for i, w in enumerate(widths):
    tbl.columns[i].width = Inches(w)
hdr = ["Nº", "ACTION", "OWNER", "DEADLINE", "STATUS"]
status_vals = [("ON TRACK", GREEN), ("ON TRACK", GREEN), ("AT RISK", YELLOW),
               ("NOT STARTED", MUTED), ("DONE", BLUE)]
for r in range(rows):
    tbl.rows[r].height = Inches(0.50 if r else 0.46)
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.08)
        cell.margin_top = cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfc = cell.text_frame
        p = tfc.paragraphs[0]
        run = p.add_run()
        f = run.font
        f.name = FONT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            run.text = hdr[c]
            f.size = Pt(9.5)
            f.bold = True
            f.color.rgb = WHITE
            _set_spc(run, 200)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else tint(NAVY, 0.035)
            if c == 0:
                run.text = f"{r}"
                f.size = Pt(10.5); f.bold = True; f.color.rgb = MUTED
            elif c == 1:
                run.text = f"[ACTION {r} — concrete, verifiable deliverable]"
                f.size = Pt(10.5); f.color.rgb = NAVY
            elif c == 2:
                run.text = "[NAME / TEAM]"
                f.size = Pt(10); f.color.rgb = BODY
            elif c == 3:
                run.text = "[DATE]"
                f.size = Pt(10); f.color.rgb = BODY
            else:
                txt_, col_ = status_vals[r - 1]
                run.text = "●  "
                f.size = Pt(9); f.bold = True; f.color.rgb = col_
                run2 = p.add_run()
                run2.text = txt_
                f2 = run2.font
                f2.name = FONT; f2.size = Pt(9); f2.bold = True
                f2.color.rgb = BODY
                _set_spc(run2, 80)
# decisions band
rect(s, 0.55, 5.72, SW - 1.1, 1.10, fill=tint(YELLOW, 0.16), rounded=True, radius=0.10)
rect(s, 0.55, 5.72, 0.09, 1.10, fill=YELLOW)
text(s, 0.88, 5.88, 4.0, 0.3, [{"runs": [
    ("DECISIONS REQUESTED FROM HQ", {"size": 9.5, "color": NAVY, "bold": True, "spc": 250})]}])
text(s, 0.88, 6.18, 11.6, 0.55, [{"runs": [
    ("[DECISION 1 — approval, budget or guideline needed]   ·   [DECISION 2]",
     {"size": 11, "color": BODY})]}])
notes(s, "PROMPTS — 1) Read only the actions that changed status; do not re-read "
         "the table. 2) Close by securing the decisions in the yellow band — get an "
         "explicit yes/no or a date from HQ before leaving the slide.")

# =========================================================== 10 · THANK YOU
s = slide_new(dark=True)
s.shapes.add_picture("build/fan_dark.png", Inches(SW - 5.3), Inches(SH - 5.3),
                     Inches(5.3), Inches(5.3))
text(s, 0.55, 1.95, 9.0, 1.3, [{"runs": [
    ("Merci. ", {"size": 50, "color": WHITE, "bold": True}),
    ("Gracias.", {"size": 50, "color": YELLOW, "bold": True}),
]}])
text(s, 0.57, 3.15, 8.0, 0.4, [{"runs": [
    ("QUESTIONS & DISCUSSION", {"size": 12, "color": LAVENDER, "spc": 380})]}])
hline(s, 0.57, 3.85, 4.2, shade_on_navy(WHITE, 0.30), 1.0)
text(s, 0.57, 4.10, 7.5, 1.3, [
    {"runs": [("[NAME]", {"size": 15, "color": WHITE, "bold": True}),
              ("  —  VP Brand Marketing, Grupo One", {"size": 12, "color": LAVENDER})],
     "space_after": 6},
    {"runs": [("[EMAIL]   ·   [PHONE]   ·   grupone.online", {"size": 11, "color": LAVENDER})]},
])
chip(s, 0.57, 5.45, 4.30, 0.50, "NEXT REVIEW: [DATE] · [FORMAT]",
     line=LAVENDER, txt_color=WHITE, size=10.5, spc=180)
# logos bottom-left
s.shapes.add_picture("build/kiabi_white.png", Inches(0.57), Inches(6.45),
                     Inches(1.55), Inches(1.55 * 320 / 1440))
gw2 = 0.62
s.shapes.add_picture("build/gone_white.png", Inches(2.55), Inches(6.32),
                     Inches(gw2), Inches(gw2 * 908 / 900))
ribbon(s)
notes(s, "PROMPTS — 1) Recap the two or three commitments made in this meeting, "
         "with owners. 2) Confirm the date and format of the next review before "
         "closing the call.")

prs.save("Kiabi_GrupoOne_Report_Template.pptx")
print("saved Kiabi_GrupoOne_Report_Template.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
